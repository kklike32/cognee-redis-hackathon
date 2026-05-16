"""Generate Sherlock hackathon pitch deck as a PowerPoint file."""
from __future__ import annotations

import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

# ── colour palette ────────────────────────────────────────────────────────────
DARK_BG   = RGBColor(0x0D, 0x0D, 0x1A)   # near-black navy
ACCENT    = RGBColor(0x5B, 0x8C, 0xFF)   # electric blue
ACCENT2   = RGBColor(0xFF, 0x6B, 0x6B)   # coral red
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY= RGBColor(0xB0, 0xB8, 0xD0)
CARD_BG   = RGBColor(0x1A, 0x1A, 0x2E)   # slightly lighter navy

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── helpers ───────────────────────────────────────────────────────────────────

def _bg(slide, color: RGBColor = DARK_BG):
    """Fill slide background with solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _box(slide, left, top, width, height,
         bg: RGBColor | None = None, border: RGBColor | None = None):
    """Add a rounded rectangle shape."""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background() if border is None else None
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if bg:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
    else:
        shape.fill.background()
    return shape


def _text_box(slide, text: str, left, top, width, height,
              font_size=20, bold=False, color: RGBColor = WHITE,
              align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def _heading(slide, text: str, top_offset=Inches(0.35)):
    """Slide section heading bar."""
    _box(slide, Inches(0), top_offset, SLIDE_W, Inches(0.6), bg=ACCENT)
    _text_box(slide, text,
              Inches(0.4), top_offset + Inches(0.07),
              SLIDE_W - Inches(0.8), Inches(0.5),
              font_size=22, bold=True, color=WHITE)


def _bullet_list(slide, items: list[str], left, top, width, height,
                 font_size=16, color: RGBColor = WHITE, bullet="•"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


def _label_value(slide, label: str, value: str, left, top, width,
                 label_size=13, value_size=15):
    """Two-line label + value pair."""
    _text_box(slide, label, left, top, width, Inches(0.3),
              font_size=label_size, color=ACCENT)
    _text_box(slide, value, left, top + Inches(0.28), width, Inches(0.4),
              font_size=value_size, color=WHITE)


def _accent_line(slide, top):
    """Thin horizontal accent rule."""
    line_shape = slide.shapes.add_shape(1, Inches(0.4), top,
                                        SLIDE_W - Inches(0.8), Pt(2))
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = ACCENT
    line_shape.line.fill.background()


# ── slides ────────────────────────────────────────────────────────────────────

def slide_title(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _bg(slide)

    # big accent stripe at top
    _box(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), bg=ACCENT)

    # logo / name block
    _text_box(slide, "🔍 SHERLOCK",
              Inches(1), Inches(1.2), Inches(11), Inches(1.4),
              font_size=72, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    _text_box(slide, "Competitive Intelligence Wiki  |  Cognee × Redis Hackathon",
              Inches(1), Inches(2.6), Inches(11), Inches(0.6),
              font_size=22, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    _accent_line(slide, Inches(3.4))

    _text_box(slide, "Self-improving battle cards for Oyster HR sales AEs",
              Inches(1), Inches(3.6), Inches(11), Inches(0.6),
              font_size=20, color=WHITE, align=PP_ALIGN.CENTER)

    # team meta row
    meta = [
        ("Team",       "Keenan"),
        ("Wiki",       "Sherlock Competitive Wiki"),
        ("Domain",     "Competitive Intelligence / Sales Enablement"),
    ]
    col_w = Inches(4)
    for i, (label, val) in enumerate(meta):
        x = Inches(0.5) + i * col_w
        _label_value(slide, label, val, x, Inches(5.0), col_w - Inches(0.2))

    _box(slide, Inches(0), Inches(7.3), SLIDE_W, Inches(0.2), bg=ACCENT2)


def slide_wiki_overview(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _heading(slide, "Wiki Overview")

    _text_box(slide,
              "Sherlock is a local-first LLM competitive-intelligence wiki for Oyster HR "
              "sales reps. It continuously ingests internal sales sources (Gong transcripts, "
              "G2 reviews, product launch docs) into a structured knowledge graph via Cognee, "
              "and caches AI-generated battle cards in Redis. A human-in-the-loop analyst "
              "review step approves proposed wiki edits, which trigger cache invalidation — "
              "so every subsequent query automatically draws from the freshest knowledge.",
              Inches(0.5), Inches(1.15), Inches(12.3), Inches(1.9),
              font_size=17, color=WHITE)

    _accent_line(slide, Inches(3.1))

    cols = [
        ("Domain / Data Sources",
         "Gong call transcripts · G2 customer reviews · competitor product-launch notes · "
         "internal battle cards"),
        ("Primary Use Case",
         "Pre-call competitive brief generation with cited evidence for Deel displacement"),
        ("What Makes It Stand Out",
         "Redis session-cache + Cognee knowledge graph + analyst HITL = a wiki that actually "
         "gets smarter after each review cycle"),
    ]
    for i, (lbl, val) in enumerate(cols):
        top = Inches(3.3) + i * Inches(1.2)
        _box(slide, Inches(0.4), top, Inches(12.5), Inches(1.1), bg=CARD_BG, border=ACCENT)
        _text_box(slide, lbl, Inches(0.6), top + Inches(0.05), Inches(12), Inches(0.3),
                  font_size=13, bold=True, color=ACCENT)
        _text_box(slide, val, Inches(0.6), top + Inches(0.35), Inches(12), Inches(0.65),
                  font_size=15, color=WHITE)


def slide_three_ops(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _heading(slide, "The Three Operations")

    ops = [
        ("⬇  INGEST",
         ACCENT,
         [
             "What goes in: Gong transcripts · G2 reviews · product-launch markdown · pasted analyst notes",
             "How captured:  File upload UI → sherlock/ingest.py → save_local_chunks() + optional cognee.add()",
             "Code entry:    sherlock/ingest.py → demo_ingest() / source_intake.py → save_source()",
         ]),
        ("🔍  QUERY + SELF-IMPROVE",
         RGBColor(0x5B, 0xFF, 0xB8),
         [
             "How users query: Streamlit Battle Card view → card_agent.py → generate_brief()",
             "Feedback source: Analyst approves / rejects proposed changes in Analyst Review UI",
             "How feedback updates wiki: approve_change() writes to wiki → invalidates Redis cache → "
             "next query retrieves enriched context",
             "Code entry:  sherlock/card_agent.py · sherlock/pending_changes.py → approve_change()",
         ]),
        ("🧹  LINT",
         ACCENT2,
         [
             "What linting means: deduplicate pending changes, resolve conflicts, prune stale proposals",
             "How it runs: on-write (every approve/reject) and on-demand from Analyst Review UI",
             "Code entry: sherlock/pending_changes.py → pending_only() · save_changes()",
         ]),
    ]

    col_w = Inches(4.2)
    for i, (title, color, bullets) in enumerate(ops):
        x = Inches(0.25) + i * (col_w + Inches(0.1))
        top = Inches(1.1)
        h = Inches(5.9)
        _box(slide, x, top, col_w, h, bg=CARD_BG, border=color)
        _text_box(slide, title, x + Inches(0.15), top + Inches(0.1),
                  col_w - Inches(0.3), Inches(0.4),
                  font_size=16, bold=True, color=color)
        _accent_line(slide, top + Inches(0.58))
        _bullet_list(slide, bullets,
                     x + Inches(0.15), top + Inches(0.65),
                     col_w - Inches(0.3), h - Inches(0.8),
                     font_size=13.5, color=WHITE, bullet="→")


def slide_self_improvement(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _heading(slide, "Self-Improvement Evidence  (Before → After)")

    # BEFORE
    _box(slide, Inches(0.3), Inches(1.1), Inches(6.1), Inches(5.8), bg=CARD_BG, border=ACCENT2)
    _text_box(slide, "BASELINE RUN", Inches(0.5), Inches(1.2), Inches(5.7), Inches(0.4),
              font_size=16, bold=True, color=ACCENT2)

    before_items = [
        ('Query', '"How does Oyster beat Deel on EOR pricing?"'),
        ('Result', 'Generic positioning — no specific pricing figures cited'),
        ('Score', '4 / 10  (judge rubric: citation depth)'),
        ('Feedback', 'error_type: missing_evidence'),
        ('', 'error_message: No pricing data found in wiki'),
        ('', 'feedback: "Brief lacks Deel price benchmarks"'),
        ('', 'success_score: 0.40'),
    ]
    top = Inches(1.75)
    for lbl, val in before_items:
        if lbl:
            _text_box(slide, lbl, Inches(0.5), top, Inches(1.5), Inches(0.32),
                      font_size=12, bold=True, color=ACCENT2)
        _text_box(slide, val, Inches(0.5 if not lbl else 2.0), top,
                  Inches(4.2 if not lbl else 4.3), Inches(0.32),
                  font_size=12, color=WHITE)
        top += Inches(0.38)

    # AFTER
    _box(slide, Inches(6.9), Inches(1.1), Inches(6.1), Inches(5.8), bg=CARD_BG, border=ACCENT)
    _text_box(slide, "IMPROVED RUN", Inches(7.1), Inches(1.2), Inches(5.7), Inches(0.4),
              font_size=16, bold=True, color=ACCENT)

    after_items = [
        ('Query', '"How does Oyster beat Deel on EOR pricing?"'),
        ('Result', 'Specific claim: "Deel charges up to 599 USD/mo vs Oyster\'s flat 499" '
                   'with [src-2] citation from G2 review'),
        ('Score', '8 / 10'),
        ('What changed', 'Analyst approved a pending change that injected a new'),
        ('', '"Quantitative Fields" section into data/wiki/deel.md'),
        ('', ''),
        ('Before wiki:', '## Quantitative Fields\\n(empty)'),
        ('After wiki:', '## Quantitative Fields\\nDeel EOR: $299–$599/seat/mo...'),
    ]
    top = Inches(1.75)
    for lbl, val in after_items:
        if lbl:
            _text_box(slide, lbl, Inches(7.1), top, Inches(1.9), Inches(0.32),
                      font_size=12, bold=True, color=ACCENT)
        _text_box(slide, val, Inches(7.1 if not lbl else 9.0), top,
                  Inches(5.7 if not lbl else 3.9), Inches(0.32),
                  font_size=12, color=WHITE)
        top += Inches(0.38)

    # arrow in the middle
    _text_box(slide, "→", Inches(6.2), Inches(3.8), Inches(0.7), Inches(0.6),
              font_size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


def slide_architecture(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _heading(slide, "Architecture")

    # left: flow diagram (text-art)
    flow = (
        "[ Sources: Gong · G2 · Docs ]\n"
        "           |\n"
        "           ▼\n"
        "[ sherlock/ingest.py ]\n"
        "           |\n"
        "    ┌──────┴───────┐\n"
        "    ▼              ▼\n"
        "[ Redis ]    [ Cognee Graph ]\n"
        "session mem   permanent KG\n"
        "    |\n"
        "    ▼\n"
        "[ card_agent.py ]\n"
        "    |\n"
        "    ▼\n"
        "[ Streamlit UI ]\n"
        "    |\n"
        "    ▼\n"
        "[ Analyst Review → approve ]\n"
        "    |\n"
        "    ▼\n"
        "[ wiki update + cache bust ]"
    )
    _box(slide, Inches(0.3), Inches(1.1), Inches(5.5), Inches(6.1), bg=CARD_BG, border=ACCENT)
    _text_box(slide, flow, Inches(0.5), Inches(1.2), Inches(5.1), Inches(5.9),
              font_size=13, color=WHITE)

    # right: Redis vs Cognee split table
    right_x = Inches(6.1)
    _text_box(slide, "Redis  —  Session Memory (Hot)",
              right_x, Inches(1.1), Inches(6.9), Inches(0.4),
              font_size=16, bold=True, color=ACCENT)
    redis_items = [
        "Stores: serialised brief JSON keyed by (competitor + deal_ctx + wiki_hash + source_hash + prompt_ver)",
        "Written: on every cache miss in card_agent.generate_brief()",
        "TTL: 24 h default, invalidated immediately on wiki approval",
        "Stays in Redis: raw response cache (per-conversation hot data)",
    ]
    _bullet_list(slide, redis_items, right_x, Inches(1.55), Inches(6.9), Inches(2.0),
                 font_size=13, color=WHITE)

    _accent_line(slide, Inches(3.7))

    _text_box(slide, "Cognee  —  Permanent Knowledge Graph",
              right_x, Inches(3.8), Inches(6.9), Inches(0.4),
              font_size=16, bold=True, color=RGBColor(0x5B, 0xFF, 0xB8))
    cognee_items = [
        "Stores: chunked markdown from all sources + wiki; entity relationships between competitors, features, pricing",
        "Promoted when: analyst approves a change → ingest.demo_ingest() re-indexes updated wiki",
        "Stays in Cognee: cross-session durable knowledge (entities, citations, relationships)",
        "Distillation improved: baseline had empty Quantitative Fields → after approval full pricing graph",
    ]
    _bullet_list(slide, cognee_items, right_x, Inches(4.25), Inches(6.9), Inches(2.5),
                 font_size=13, color=WHITE)


def slide_agents(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _heading(slide, "Agents & Roles")

    roles = [
        ("Ingestor",  ACCENT,
         "sherlock/ingest.py\nsherlock/source_intake.py\n\n"
         "Reads .md/.txt/.html/.pdf → normalises → saves to data/sources/ → "
         "calls save_local_chunks() + optional Cognee add()"),
        ("Querier / Card Agent", RGBColor(0x5B, 0xFF, 0xB8),
         "sherlock/card_agent.py\n\n"
         "Builds cache key → Redis lookup → on miss: retrieve_context_with_status() "
         "→ build_citations() → LLM or deterministic brief → Redis set()"),
        ("Linter / Critic", ACCENT2,
         "sherlock/pending_changes.py\n\n"
         "pending_only() surfaces unreviewed proposals; approve_change() / reject_change() "
         "applies edits to wiki, prunes stale pending JSON, invalidates Redis"),
        ("Wiki Builder",  RGBColor(0xFF, 0xD7, 0x00),
         "sherlock/wiki_builder.py\n\n"
         "build_company_wiki() extracts structured sections from sources; "
         "optional LLM synthesis gated by SHERLOCK_USE_LLM env; "
         "deterministic fallback always available"),
    ]

    col_w = Inches(3.15)
    for i, (name, color, desc) in enumerate(roles):
        x = Inches(0.2) + i * (col_w + Inches(0.1))
        _box(slide, x, Inches(1.1), col_w, Inches(5.8), bg=CARD_BG, border=color)
        _text_box(slide, name, x + Inches(0.12), Inches(1.18), col_w - Inches(0.25), Inches(0.4),
                  font_size=15, bold=True, color=color)
        _accent_line(slide, Inches(1.65))
        _text_box(slide, desc, x + Inches(0.12), Inches(1.72), col_w - Inches(0.25), Inches(5.1),
                  font_size=12.5, color=WHITE)


def slide_reproduction(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _heading(slide, "Reproduction & Environment")

    code = (
        "# 1. Clone & install\n"
        "git clone <repo-url> && cd cognee-redis-hackathon\n"
        "pip install -e .\n\n"
        "# 2. Set environment variables\n"
        "export OPENAI_API_KEY=sk-...\n"
        "export REDIS_URL=redis://localhost:6379\n"
        "export SHERLOCK_USE_LLM=1          # optional\n\n"
        "# 3. Start Redis (Docker)\n"
        "docker compose up -d redis\n\n"
        "# 4. Ingest demo data\n"
        "python scripts/ingest_demo_data.py\n\n"
        "# 5. Launch UI\n"
        "streamlit run app/streamlit_app.py\n\n"
        "# 6. Smoke test\n"
        "python scripts/smoke_test.py"
    )
    _box(slide, Inches(0.3), Inches(1.1), Inches(7.8), Inches(6.1), bg=CARD_BG, border=ACCENT)
    _text_box(slide, code, Inches(0.5), Inches(1.2), Inches(7.5), Inches(5.9),
              font_size=13, color=RGBColor(0xA8, 0xFF, 0xA8))

    env_items = [
        "OPENAI_API_KEY   — LLM calls (optional, falls back deterministically)",
        "REDIS_URL        — Redis connection (default: redis://localhost:6379)",
        "SHERLOCK_USE_LLM — enable GPT synthesis in wiki builder",
        "SHERLOCK_INDEX_COGNEE — enable Cognee ingestion",
        "COGNEE_LLM_API_KEY   — if different from OPENAI_API_KEY",
    ]
    _text_box(slide, "Required Environment Variables",
              Inches(8.4), Inches(1.1), Inches(4.7), Inches(0.5),
              font_size=15, bold=True, color=ACCENT)
    _bullet_list(slide, env_items, Inches(8.4), Inches(1.65), Inches(4.7), Inches(3.0),
                 font_size=12.5, color=WHITE)

    _text_box(slide, "Tests",
              Inches(8.4), Inches(4.8), Inches(4.7), Inches(0.4),
              font_size=15, bold=True, color=ACCENT)
    _bullet_list(slide,
                 ["pytest tests/ — full unit + integration suite",
                  "9 test modules covering cache, card agent, citations,",
                  "  markdown store, retrieval, HITL, source intake, wiki builder"],
                 Inches(8.4), Inches(5.25), Inches(4.7), Inches(1.5),
                 font_size=12.5, color=WHITE)


def slide_demo_outline(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _heading(slide, "3-Minute Pitch Outline & Demo")

    steps = [
        ("1  Problem / Idea",
         "Sales AEs waste 20+ min hunting for competitive intel before every call. "
         "Sherlock delivers a cited, deal-specific battle card in under 3 seconds."),
        ("2  Ingest Demo",
         "Upload a Gong transcript → Source Intake UI → file saved to data/sources/ → "
         "Cognee indexes entities → wiki rebuild triggered"),
        ("3  Query Demo (Before)",
         "AE enters deal context → card_agent generates brief → Redis MISS → LLM call → "
         "brief stored. Note: generic pricing section (no figures)"),
        ("4  Self-Improve Step",
         "Analyst Review UI → approve pending change with pricing data → "
         "wiki updated → Redis cache invalidated for Deel key-space"),
        ("5  Query Demo (After)",
         "Same query → Redis HIT latency ~12 ms vs ~1400 ms → brief now includes "
         "specific pricing benchmarks with [src-2] citation"),
        ("6  What's Next",
         "Multi-competitor support · Slack ingest connector · Auto-linting with GPT-4o critic "
         "· Confidence scoring on citations · Cloud deployment"),
    ]

    col_h = Inches(1.0)
    for i, (title, desc) in enumerate(steps):
        row = i // 2
        col = i % 2
        x = Inches(0.3) + col * Inches(6.55)
        top = Inches(1.15) + row * (col_h + Inches(0.12))
        _box(slide, x, top, Inches(6.4), col_h, bg=CARD_BG, border=ACCENT)
        _text_box(slide, title, x + Inches(0.12), top + Inches(0.05),
                  Inches(6.1), Inches(0.3), font_size=13, bold=True, color=ACCENT)
        _text_box(slide, desc, x + Inches(0.12), top + Inches(0.35),
                  Inches(6.1), Inches(0.58), font_size=12, color=WHITE)

    _text_box(slide, "Links",
              Inches(0.3), Inches(5.5), Inches(12.7), Inches(0.35),
              font_size=15, bold=True, color=ACCENT)
    _bullet_list(slide,
                 ["Repo:   github.com/keenan/cognee-redis-hackathon",
                  "Slides: this deck",
                  "Demo:   Loom / YouTube link (to be added)"],
                 Inches(0.3), Inches(5.9), Inches(12.7), Inches(1.4),
                 font_size=13, color=WHITE)


# ── main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_wiki_overview(prs)
    slide_three_ops(prs)
    slide_self_improvement(prs)
    slide_architecture(prs)
    slide_agents(prs)
    slide_reproduction(prs)
    slide_demo_outline(prs)

    out = Path(__file__).parent.parent / "sherlock_pitch.pptx"
    prs.save(str(out))
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
